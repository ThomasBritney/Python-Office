# 书写人：胡高远
# 书写时间：2022/7/21  15:39
from pptx import Presentation
pre = Presentation('算法复杂性问题.pptx')
slides = pre.slides     #获取PPT所有幻灯片
for slide in slides:
    # print(slide)
    for shape in slide.shapes:
        print(shape)         #获取shape形状
    print('---------------------')







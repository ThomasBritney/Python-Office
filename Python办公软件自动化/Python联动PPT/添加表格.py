# 书写人：胡高远
# 书写时间：2022/7/23  10:15
from pptx import Presentation
from pptx.util import Cm
# 新建PPT文件
prs = Presentation()

# 添加一张幻灯片  （应用6号幻灯片）
slide = prs.slides.add_slide(prs.slide_layouts[6])

rows,cols = 4,2

left = top = Cm(5)
width = Cm(8)
height = Cm(4)
table = slide.shapes.add_table(rows,cols,left,top,width,height).table

# 准备数据
lst = [
    ['学号','姓名'],
    ['1001','张三'],
    ['1002','李四'],
    ['1003','王二']
]
for row in range(rows):
    for col in range(cols):
        table.cell(row,col).text = str(lst[row][col])


prs.save('添加表格.pptx')
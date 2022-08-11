# 书写人：胡高远
# 书写时间：2022/7/22  16:15
from pptx import Presentation
from docx import Document

pre = Presentation('算法复杂性问题.pptx')
lst = []       #用于存储每一个段落
for slide in pre.slides:
    #获取shape形状
    for shape in slide.shapes:
        #判断形状是否存在文字
        if shape.has_text_frame:
            #获取文字框
            text_frame = shape.text_frame
            for para in text_frame.paragraphs:
                lst.append(para.text)
# print(lst)
doc = Document()
for item in lst:
    doc.add_paragraph(item)
doc.save('算法复杂性问题.docx')




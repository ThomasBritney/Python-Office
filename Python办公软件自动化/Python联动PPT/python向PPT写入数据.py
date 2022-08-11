# 书写人：胡高远
# 书写时间：2022/7/22  16:35
from pptx import Presentation
#  新建PPT文件
prs = Presentation()
# 添加一张幻灯片
slide = prs.slides.add_slide(prs.slide_layouts[0])
for shape in slide.placeholders:
    phf = shape.placeholder_format
    print(phf.idx)
    print(shape.name)
    print(phf.type)
    print('----------------')




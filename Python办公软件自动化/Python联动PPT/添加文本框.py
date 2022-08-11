# 书写人：胡高远
# 书写时间：2022/7/23  9:45
from pptx import Presentation
from pptx.util import Cm
# 新建PPT文件
prs = Presentation()

# 添加一张幻灯片  （应用6号幻灯片）
slide = prs.slides.add_slide(prs.slide_layouts[6])

left=top=width=height=Cm(3)
#添加文本框
text_box = slide.shapes.add_textbox(left,top,width,height)
tf = text_box.text_frame
tf.text = '小荷才露尖尖角，早有蜻蜓立上头'

# 保存成PPT文件
prs.save('添加文本框.pptx')
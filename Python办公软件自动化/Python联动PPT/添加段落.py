# 书写人：胡高远
# 书写时间：2022/7/23  9:45
from pptx import Presentation
# 新建PPT文件
prs = Presentation()
# 添加一张幻灯片  （应用一号幻灯片）
slide = prs.slides.add_slide(prs.slide_layouts[1])
# 获取形状
shapes = slide.shapes
# 标题与正文
title_shape = shapes.title     #获取标题的形状
body_shape = shapes.placeholders[1]    #获取正文的占位符

title_shape.text = '使用Python向PPT文件中写入内容'

# 获取正文的文本框
tf = body_shape.text_frame   #文本框本身也是一个段落
tf.text = '新建PPT文件'
# 在文本框中添加段落
p = tf.add_paragraph()
p.text = '添加新的幻灯片'
p.level = 1            #添加层级
p = tf.add_paragraph()
p.text = '保存操作'
p.level = 1             #添加层级
# 保存成PPT文件
prs.save('添加段落.pptx')
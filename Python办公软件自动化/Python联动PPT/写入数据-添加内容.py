# 书写人：胡高远
# 书写时间：2022/7/22  16:44
from pptx import Presentation
#  新建PPT文件
prs = Presentation()
# 添加一张幻灯片
slide = prs.slides.add_slide(prs.slide_layouts[0])

shapes = slide.placeholders
shapes[0].text = '使用Python向PPT写入数据'
shapes[1].text = '-----胡高远'

# 保存
prs.save('添加内容.pptx')



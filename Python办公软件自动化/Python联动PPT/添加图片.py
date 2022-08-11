# 书写人：胡高远
# 书写时间：2022/7/23  10:15
from pptx import Presentation
from pptx.util import Cm
# 新建PPT文件
prs = Presentation()

# 添加一张幻灯片  （应用6号幻灯片）
slide = prs.slides.add_slide(prs.slide_layouts[6])

left=top=width=height=Cm(5)
pic = slide.shapes.add_picture('logo.png',left,top,width,height)

left = Cm(12)
top = Cm(6)
# width = Cm(5)
height = Cm(9)
pic = slide.shapes.add_picture('logo.png',left,top,height=height)


prs.save('添加图片.pptx')
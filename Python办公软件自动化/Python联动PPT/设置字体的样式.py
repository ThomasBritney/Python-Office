# 书写人：胡高远
# 书写时间：2022/7/24  15:35
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import MSO_VERTICAL_ANCHOR,PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from docx.shared import Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 添加文本框
left=top=height=Cm(3)
width = Cm(10)
textbox = slide.shapes.add_textbox(left,top,width,height)
# 向文本框中添加段落
tf = textbox.text_frame
p = tf.add_paragraph()
p.text = '物是人非事事休，欲语泪先流'
p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
# 设置段落字体样式
p.font.name = '微软雅黑'
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(200,200,200)

# 设置样式
tf.margin_bottom = Cm(1)
tf.margin_top = Cm(1)
tf.margin_left = Cm(1)
tf.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
tf.word_wrap = True

# 设置文本框的背景颜色
fill = textbox.fill
fill.solid()   #纯色填充
fill.fore_color.rgb = RGBColor(200,200,100)

# 文本框的边框样式
line = textbox.line
line.color.rgb = RGBColor(255,0,0)
line.width = Cm(0.2)

prs.save('设置字体样式.pptx')